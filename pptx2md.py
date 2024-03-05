import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPEº
from pptx.enum.text import PP_ALIGN

def convert_pptx_to_md(pptx_dir):
    # Lista todos los archivos en el directorio
    for filename in os.listdir(pptx_dir):
        # Procesar solo los archivos .pptx
        if filename.endswith(".pptx"):
            file_path = os.path.join(pptx_dir, filename)
            markdown_text = ""

            # Leer la presentación
            presentation = Presentation(file_path)
            
            # Extraer el texto de la portada
            first_slide = presentation.slides[0]
            title = first_slide.shapes.title.text.strip()
            markdown_text += f"# {title}\n\n"
            
            # Iterar sobre las diapositivas y extraer el contenido
            for slide in presentation.slides:
                # Extraer el título principal de la diapositiva
                title_shape = slide.shapes.title
                if title_shape:
                    main_title = title_shape.text.strip()
                    markdown_text += f"## {main_title}\n\n"
                
                # Extraer el texto de las demás formas
                for shape in slide.shapes:
                    # Ignorar los grupos de formas (por su complejidad)
                    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                        continue
                    if shape.has_text_frame:
                        text_frame = shape.text_frame
                        for paragraph in text_frame.paragraphs:
                            if paragraph.text.strip():
                                # Agregar el texto como un bullet point
                                markdown_text += f"- {paragraph.text.strip()}\n"
            
            # Guardar el texto en un archivo Markdown
            markdown_filename = os.path.splitext(filename)[0] + ".md"
            markdown_path = os.path.join(pptx_dir, markdown_filename)
            with open(markdown_path, "w") as md_file:
                md_file.write(markdown_text)

# Directorio que contiene los archivos .pptx
pptx_directory = "/content/sample_data"
convert_pptx_to_md(pptx_directory)
